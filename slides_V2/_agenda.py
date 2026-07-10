# -*- coding: utf-8 -*-
"""Rebuild slide 2 (Agenda) as the PDF/beamer style: blue numbered spheres +
blue text labels on white, vertical list. No cards."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

BLUE_TXT = RGBColor(0x1F, 0x60, 0x91)   # label blue (medium)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SPH_IN = "3D7AB5"   # sphere highlight (center)
SPH_OUT = "0E2C49"  # sphere edge (dark navy)

ITEMS = ["Introduction", "Perceiver", "Experiments: Perceiver",
         "Perceiver IO", "Experiments: Perceiver IO", "Conclusion"]

def radial_sphere(shape, inner, outer):
    """Replace the shape fill with a radial (circle) gradient -> 3D ball look."""
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(inner)
    spPr = shape._element.spPr
    old = spPr.find(qn('a:solidFill'))
    grad = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, clr in [(0, inner), (100000, outer)]:
        gs = etree.SubElement(gsLst, qn('a:gs')); gs.set('pos', str(pos))
        etree.SubElement(gs, qn('a:srgbClr')).set('val', clr)
    path = etree.SubElement(grad, qn('a:path')); path.set('path', 'circle')
    fr = etree.SubElement(path, qn('a:fillToRect'))
    fr.set('l', '30000'); fr.set('t', '20000'); fr.set('r', '70000'); fr.set('b', '80000')
    # move gradFill where solidFill was (before a:ln), then drop solidFill
    if old is not None:
        old.addprevious(grad); spPr.remove(old)
    shape.line.fill.background()
    shape.shadow.inherit = False

p = Presentation("Perceiver & Perceiver IO.pptx")
s = p.slides[1]

# delete every shape except footer [0] and title [1]
for sh in list(s.shapes)[2:]:
    sh._element.getparent().remove(sh._element)

top0, pitch, D = 2.2, 0.82, 0.62
for i, name in enumerate(ITEMS):
    cy = top0 + i * pitch
    # sphere
    ball = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(cy), Inches(D), Inches(D))
    radial_sphere(ball, SPH_IN, SPH_OUT)
    tf = ball.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    run = pr.add_run(); run.text = str(i + 1)
    run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = WHITE
    # label
    lbl = s.shapes.add_textbox(Inches(1.85), Inches(cy - 0.08), Inches(10.6), Inches(D + 0.16))
    lt = lbl.text_frame; lt.word_wrap = True; lt.vertical_anchor = MSO_ANCHOR.MIDDLE
    lp = lt.paragraphs[0]
    lr = lp.add_run(); lr.text = name
    lr.font.size = Pt(27); lr.font.color.rgb = BLUE_TXT

p.save("Perceiver & Perceiver IO.pptx")
print("Agenda rebuilt as PDF-style numbered spheres")
